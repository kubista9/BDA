import os
import subprocess
from pptx import Presentation

# Directory containing your files
dir_path = '/Users/jakubkuka/Desktop/BDA/sources'

# Extract text from PDF files using pdftotext
for filename in os.listdir(dir_path):
    file_path = os.path.join(dir_path, filename)
    
    # For PDF files
    if filename.endswith('.pdf'):
        output_txt = f"{file_path}.txt"
        subprocess.run(['pdftotext', file_path, output_txt])
        print(f"Text extracted from: {filename}")
    
    # For PowerPoint files
    if filename.endswith('.pptx'):
        pptx_txt = f"{file_path}.txt"
        prs = Presentation(file_path)
        with open(pptx_txt, 'w') as txt_file:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        txt_file.write(shape.text + "\n")
        print(f"Text extracted from: {filename}")
